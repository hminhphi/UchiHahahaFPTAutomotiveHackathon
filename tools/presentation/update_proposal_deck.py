from __future__ import annotations

from pathlib import Path
from shutil import copyfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[2]
PPTX_PATH = ROOT / "Template for Teams - Hackathon 2026.pptx"
BACKUP_PATH = ROOT / "artifacts" / "Template for Teams - Hackathon 2026.backup.pptx"
ASSET_DIR = ROOT / "artifacts" / "pptx-assets"

NAVY = RGBColor(0x19, 0x22, 0x6D)
ORANGE = RGBColor(0xF3, 0x70, 0x21)
BLUE = RGBColor(0x03, 0x4E, 0xA2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x7A, 0x90)


def set_text(shape, text: str, *, size: int | None = None, bold: bool | None = None, color: RGBColor | None = None, align: PP_ALIGN | None = None) -> None:
    shape.text = text
    if not hasattr(shape, "text_frame"):
        return
    shape.text_frame.word_wrap = True
    for paragraph in shape.text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            if size is not None:
                run.font.size = Pt(size)
            if bold is not None:
                run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color
            run.font.name = "Quattrocento Sans"


def format_text_frame(shape, size: int, *, bold: bool = False, color: RGBColor = DARK, align: PP_ALIGN | None = None) -> None:
    if not hasattr(shape, "text_frame"):
        return
    shape.text_frame.word_wrap = True
    for paragraph in shape.text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Quattrocento Sans"


def cover_picture(slide, shape, image_name: str) -> None:
    slide.shapes.add_picture(str(ASSET_DIR / image_name), shape.left, shape.top, width=shape.width, height=shape.height)


def normalize_footer(shape, text: str) -> None:
    set_text(shape, text, size=9, color=WHITE)


def main() -> None:
    BACKUP_PATH.parent.mkdir(parents=True, exist_ok=True)
    copyfile(PPTX_PATH, BACKUP_PATH)

    prs = Presentation(str(PPTX_PATH))

    # Slide 1
    slide = prs.slides[0]
    set_text(slide.shapes[1], "FLEETIQ GUARDIAN\nHACKATHON PROPOSAL", size=28, bold=True, color=WHITE)
    set_text(slide.shapes[2], "Hanoi, Jul 10th 2026", size=11, color=WHITE)
    set_text(slide.shapes[4], "Challenge #3 | 4 AI + 1 Automotive C++", size=20, color=WHITE)

    # Slide 2
    slide = prs.slides[1]
    set_text(slide.shapes[0], "VÌ SAO PROPOSAL NÀY CÓ THỂ ĐI TIẾP", size=26, bold=True, color=NAVY)
    table = slide.shapes[1].table
    table_data = [
        ("Tiêu chí", "Điểm"),
        ("Ý tưởng", "Driver Intelligence Platform giải đúng bài toán giám sát đội xe từ xa, có giá trị vận hành rõ ràng.", "35"),
        ("Tính khả thi", "2-3 tuần đủ để hoàn thành MVP theo hướng baseline-plus: TTC, scoring, fusion, dashboard.", "30"),
        ("Hiểu đề & starter pack", "Tận dụng đủ road camera, driver camera, depth, calibration, labels, telemetry và baseline TTC.", "20"),
        ("Năng lực đội", "4 AI members xử lý perception/fusion/dashboard data, 1 Automotive C++ member lo simulator bridge và dữ liệu xe.", "15"),
    ]
    # header
    for col in range(2):
        table.cell(0, col).text = table_data[0][col]
    rows = table_data[1:]
    for row_idx, (title, body, score) in enumerate(rows, start=1):
        table.cell(row_idx, 0).text = f"{title}: {body}"
        table.cell(row_idx, 1).text = score
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = "Quattrocento Sans"
                    run.font.size = Pt(14 if r == 0 else 12)
                    run.font.bold = r == 0
                    run.font.color.rgb = DARK if r > 0 else WHITE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif c == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF7, 0xE3, 0xD8)

    # Slide 3
    slide = prs.slides[2]
    set_text(slide.shapes[0], "NỘI DUNG TRÌNH BÀY", size=24, bold=True, color=WHITE)
    toc = [
        ("1", "Đội hình & cách chia việc"),
        ("2", "Bài toán chọn thi & vì sao"),
        ("3", "Giải pháp FleetIQ Guardian"),
        ("4", "Kiến trúc, roadmap, impact"),
    ]
    for idx, (num, label) in enumerate(toc):
        set_text(slide.shapes[1 + idx * 2], num, size=26, bold=True, color=WHITE)
        set_text(slide.shapes[2 + idx * 2], label, size=20, bold=True, color=WHITE)

    # Slide 4
    slide = prs.slides[3]
    set_text(slide.shapes[0], "ĐỘI HÌNH THỰC THI", size=26, bold=True, color=NAVY)
    set_text(slide.shapes[1], "04", size=10, color=MUTED)
    set_text(slide.shapes[2], "Tên đội chơi: FleetIQ Guardian | 5 thành viên", size=18, bold=True, color=NAVY)
    cards = [
        ("A", "Perception & TTC", "2 AI members", "Stereo/depth, object tracking, TTC stream, near-miss detection."),
        ("B", "Driver State & Fusion", "1 AI member", "Driver state timeline, compound-risk logic, coaching rules."),
        ("C", "Backend + Data Layer", "1 AI member", "Event schema, scoring API, static JSON, dashboard contract."),
        ("D", "Automotive C++ Bridge", "1 member", "Simulator bridge, telemetry parser, signal validation, demo stability."),
    ]
    card_indices = [(5, 6, 7, 8), (11, 12, 13, 14), (17, 18, 19, 20), (23, 24, 25, 26)]
    for (letter_idx, name_idx, role_idx, body_idx), (letter, name, role, body) in zip(card_indices, cards):
        set_text(slide.shapes[letter_idx], letter, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        set_text(slide.shapes[name_idx], name, size=18, bold=True, color=NAVY)
        set_text(slide.shapes[role_idx], role, size=14, bold=True, color=ORANGE)
        set_text(slide.shapes[body_idx], body, size=12, color=DARK)

    # Slide 5
    slide = prs.slides[4]
    set_text(slide.shapes[0], "NĂNG LỰC ĐỘI & HÌNH DUNG SẢN PHẨM", size=26, bold=True, color=NAVY)
    set_text(slide.shapes[1], "05", size=10, color=MUTED)
    set_text(slide.shapes[2], "Cơ cấu 4 AI + 1 Automotive C++", size=18, bold=True, color=NAVY)
    set_text(slide.shapes[3], "Mockup dashboard để chốt câu chuyện demo", size=18, bold=True, color=NAVY)
    cover_picture(slide, slide.shapes[4], "team-composition.png")
    cover_picture(slide, slide.shapes[5], "dashboard-preview.png")

    # Slide 6
    slide = prs.slides[5]
    set_text(slide.shapes[0], "BÀI TẬP LỰA CHỌN", size=26, bold=True, color=NAVY)
    set_text(slide.shapes[1], "Challenge #3 làm bài nộp, Challenge #1 và #2 là engine lõi", size=21, bold=True, color=NAVY)
    headings = [
        "Tín hiệu đầu vào",
        "Engine AI",
        "Đầu ra sản phẩm",
        "MVP 2-3 tuần",
    ]
    bodies = [
        "Road camera đa góc, driver camera, depth, calibration, labels, telemetry sensor-fusion.",
        "Trip scoring, TTC monitor, near-miss eventing, driver-state fusion, confidence logic.",
        "Fleet dashboard, driver ranking, trip detail, event evidence, coaching report.",
        "Một trip chạy end-to-end với timeline đồng bộ, score, TTC, evidence và fallback static JSON.",
    ]
    for idx in range(4):
        set_text(slide.shapes[2 + idx * 2], headings[idx], size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        set_text(slide.shapes[3 + idx * 2], bodies[idx], size=12, color=DARK, align=PP_ALIGN.CENTER)
    normalize_footer(slide.shapes[14], "FleetIQ Guardian | 06")
    for pic_idx, image_name in zip([10, 11, 12, 13], ["icon-road.png", "icon-driver.png", "icon-fusion.png", "icon-dashboard.png"]):
        cover_picture(slide, slide.shapes[pic_idx], image_name)

    # Slide 7
    slide = prs.slides[6]
    set_text(slide.shapes[0], "VẤN ĐỀ CẦN GIẢI", size=26, bold=True, color=NAVY)
    set_text(slide.shapes[1], "Fleet Manager có dữ liệu nhưng chưa có quyết định rủi ro có thể hành động", size=22, bold=True, color=NAVY)
    set_text(slide.shapes[3], "6", size=44, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    set_text(slide.shapes[4], "nguồn tín hiệu cần hợp nhất thành một timeline duy nhất để giải thích vì sao chuyến xe rủi ro", size=13, color=DARK, align=PP_ALIGN.CENTER)
    issue_texts = [
        ("Dữ liệu rời rạc", "Camera, depth, telemetry và nhãn baseline đang nằm ở các luồng khác nhau; không ghép lại thì khó kể câu chuyện."),
        ("TTC thô dễ nhiễu", "Chỉ nhìn từng frame sẽ sinh false alarm; proposal cần smoothing, confidence và event merging."),
        ("Thiếu bằng chứng hành động", "Người quản lý đội xe cần biết ai nguy hiểm, vì sao, ở thời điểm nào và nên coaching điều gì."),
    ]
    for base_idx, (title, body) in zip([7, 11, 15], issue_texts):
        set_text(slide.shapes[base_idx], title, size=18, bold=True, color=NAVY)
        set_text(slide.shapes[base_idx + 1], body, size=12, color=DARK)

    # Slide 8
    slide = prs.slides[7]
    set_text(slide.shapes[0], "Ý TƯỞNG TRIỂN KHAI", size=26, bold=True, color=NAVY)
    set_text(slide.shapes[1], "FleetIQ Guardian biến camera + telemetry thành bằng chứng rủi ro cho đội xe", size=21, bold=True, color=NAVY)
    set_text(
        slide.shapes[3],
        "Elevator pitch: Chúng tôi xây một lớp intelligence cho Fleet Manager, nơi road camera, driver camera, depth và telemetry được hợp nhất thành trip score, short-TTC alerts, near-miss evidence và coaching recommendation thay vì chỉ là log kỹ thuật rời rạc.",
        size=13,
        color=DARK,
    )
    set_text(slide.shapes[6], "Insight cốt lõi", size=18, bold=True, color=NAVY)
    set_text(
        slide.shapes[7],
        "Điểm mạnh không nằm ở một model đơn lẻ. Điểm mạnh nằm ở việc nối Challenge #1 scoring và Challenge #2 TTC thành Challenge #3: một dashboard mà judge có thể hiểu trong 30 giây, drill-down trong 1 phút và tin rằng team làm xong trong 2-3 tuần.",
        size=13,
        color=DARK,
    )
    cover_picture(slide, slide.shapes[5], "icon-insight.png")

    # Slide 9
    slide = prs.slides[8]
    set_text(slide.shapes[0], "KIẾN TRÚC GIẢI PHÁP", size=26, bold=True, color=NAVY)
    normalize_footer(slide.shapes[1], "FleetIQ Guardian | 09")
    column_data = [
        ("Nguồn dữ liệu", "Road camera\ndriver camera\ndepth + calib\ntelemetry + labels", "Input"),
        ("Core engines", "TTC stream\ndriver state\ntelemetry events\ntrip scoring", "AI"),
        ("Intelligence layer", "event schema\nconfidence\ncompound risk\nexplanation", "Fusion"),
        ("Outputs", "fleet dashboard\ntrip drill-down\nevidence panel\ncoaching report", "Product"),
    ]
    shape_sets = [(15, 16, 18), (20, 21, 23), (25, 26, 28), (30, 31, 33)]
    for (title_idx, body_idx, tag_idx), (title, body, tag) in zip(shape_sets, column_data):
        set_text(slide.shapes[title_idx], title, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        set_text(slide.shapes[body_idx], body, size=12, color=DARK, align=PP_ALIGN.CENTER)
        set_text(slide.shapes[tag_idx], tag, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 10
    slide = prs.slides[9]
    set_text(slide.shapes[0], "ROADMAP TRIỂN KHAI 2-3 TUẦN", size=26, bold=True, color=NAVY)
    normalize_footer(slide.shapes[1], "FleetIQ Guardian | 10")
    roadmap = [
        ("Tuần 1", "Dựng data contract", "Load được 1 trip, chuẩn hoá event schema, đọc baseline TTC, dựng static JSON cho dashboard."),
        ("Tuần 2", "Hoàn thành 2 engine", "Trip scoring + TTC smoothing, near-miss event windows, driver-state segments, confidence logic."),
        ("Tuần 3", "Fusion & demo", "Trip drill-down, compound risk, coaching recommendation, slide/demo/video backup nếu live model lỗi."),
    ]
    shape_sets = [(4, 5, 6), (9, 10, 11), (14, 15, 16)]
    for (week_idx, title_idx, body_idx), (week, title, body) in zip(shape_sets, roadmap):
        set_text(slide.shapes[week_idx], week, size=18, bold=True, color=WHITE)
        set_text(slide.shapes[title_idx], title, size=18, bold=True, color=NAVY)
        set_text(slide.shapes[body_idx], body, size=12, color=DARK)
    set_text(slide.shapes[17], "■ Must-have: score + TTC + one fused trip | ■ Nice-to-have: annotated video export, route heatmap, live alerts", size=11, color=MUTED)

    # Slide 11
    slide = prs.slides[10]
    set_text(slide.shapes[0], "TẦM NHÌN SAU HACKATHON", size=26, bold=True, color=NAVY)
    set_text(slide.shapes[1], "Không chỉ là demo AI; đây là hạt nhân của remote fleet safety operations", size=22, bold=True, color=NAVY)
    set_text(slide.shapes[2], "Phục vụ ai", size=17, bold=True, color=NAVY)
    set_text(slide.shapes[3], "Fleet Manager, OEM analytics team và đơn vị vận hành đội xe cần xếp hạng rủi ro và coaching có bằng chứng.", size=13, color=DARK)
    set_text(slide.shapes[4], "Vì sao đáng làm tiếp", size=17, bold=True, color=NAVY)
    set_text(slide.shapes[5], "Kiến trúc tách input, engine, intelligence, output nên có thể mở rộng sang nhiều xe, nhiều loại cảnh báo và báo cáo hậu kiểm.", size=13, color=DARK)
    set_text(slide.shapes[6], "Bước tiếp theo", size=17, bold=True, color=NAVY)
    set_text(slide.shapes[7], "Sau vòng proposal, team sẽ ưu tiên dashboard thật, evidence clip có overlay TTC, và benchmark qualitative so với baseline starter pack.", size=13, color=DARK)
    set_text(slide.shapes[8], "Nguyên tắc xuyên suốt: explainable trước, ổn định trước, rồi mới thêm mô hình nặng hơn.", size=12, color=MUTED)

    # Slide 12
    slide = prs.slides[11]
    set_text(slide.shapes[0], "THANK YOU!", size=30, bold=True, color=WHITE)
    set_text(
        slide.shapes[1],
        "FleetIQ Guardian\nChallenge #3 proposal built for a 5-member team\n4 AI members + 1 Automotive C++ member | 2-3 week execution plan",
        size=16,
        color=WHITE,
    )
    format_text_frame(slide.shapes[1], 16, color=WHITE)

    prs.save(str(PPTX_PATH))


if __name__ == "__main__":
    main()
